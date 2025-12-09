# =======================================================================
# i3T4AN (Ethan Blair)
# Project:      Vector Knowledge Base
# File:         PowerPoint document text extractor
# =======================================================================

from typing import Tuple, Dict, Any
from pptx import Presentation
import logging
from .base import BaseExtractor
from exceptions import ExtractionError

logger = logging.getLogger(__name__)

class PptxExtractor(BaseExtractor):
    def extract(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from PowerPoint presentation.
        
        Args:
            file_path: Path to the .pptx file
            
        Returns:
            Tuple containing:
            - Extracted text content
            - Metadata dictionary
        """
        try:
            text_content = []
            metadata = {"slide_count": 0}
            
            prs = Presentation(file_path)
            metadata["slide_count"] = len(prs.slides)
            
            # Extract core properties if available
            if prs.core_properties:
                if prs.core_properties.title:
                    metadata["title"] = prs.core_properties.title
                if prs.core_properties.author:
                    metadata["author"] = prs.core_properties.author
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    # Join text within a slide with newlines
                    slide_content = "\n".join(slide_text)
                    text_content.append(slide_content)
            
            # Join slides with double newlines to separate them clearly
            return "\n\n".join(text_content), metadata
            
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {str(e)}")
            raise ExtractionError(f"Failed to extract PPTX: {str(e)}")
